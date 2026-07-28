import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_v3_chapter1_pptx():
    prs = Presentation()
    # 16:9 Aspect Ratio (13.333 in x 7.5 in)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Tokens v3 (Forensic Noir v3)
    COLOR_BG_SLATE = RGBColor(13, 18, 29)     # #0D121D
    COLOR_BG_NAVY = RGBColor(8, 12, 20)       # #080C14
    COLOR_CARD_BG = RGBColor(21, 29, 42)      # #151D2A
    COLOR_GOLD = RGBColor(255, 183, 51)       # #FFB733 Accent Gold
    COLOR_RED = RGBColor(255, 77, 79)         # #FF4D4F Signal Crimson
    COLOR_TEAL = RGBColor(54, 207, 201)       # #36CFC9 Vibrant Teal
    COLOR_TEXT_MAIN = RGBColor(245, 247, 250) # #F5F7FA Off-White
    COLOR_TEXT_SUB = RGBColor(196, 209, 223)  # #C4D1DF Soft Ice Blue

    def set_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Pretendard"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_SUB

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # Slide Data List v3 (p1 ~ p12)
    slides_data = [
        # p1: Cover [T1: Hero Headline]
        {
            "bg": COLOR_BG_NAVY,
            "notes": "『안녕하십니까. 충남 논산시 소재의, 노성중 교사 최재석입니다. 오늘이 1정 연수 마지막주차라고 들었는데 얼마나 선생님들이 지친 상태일지 강사인 제가 걱정이 됩니다. 힘들었다고 대답하시더라구요. 제가 덜 힘들게 해드릴 순 없고, 편안하게 들으실 수 있게, 그리고 한가지라도 유익한 정보를 남길 수 있도록 최선에 다해보겠습니다.』",
            "build": lambda s: [
                set_background(s, COLOR_BG_NAVY),
                s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(1.2), Inches(4.1), Inches(0.55)).fill.solid(),
                setattr(s.shapes[-1].fill.fore_color, 'rgb', COLOR_GOLD),
                setattr(s.shapes[-1].line.color, 'rgb', COLOR_GOLD),
                # Title
                s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5)).text_frame.paragraphs[0].line_spacing_set(1.2) or
                s.shapes[-1].text_frame.add_paragraph().line_spacing_set(1.2) or None,
            ]
        }
    ]

    # Rebuilding clean slide-by-slide pptx structure for 12 slides
    slide_specs = [
        # 1: Cover
        ("p1", COLOR_BG_NAVY, "『안녕하십니까. 충남 논산시 소재의, 노성중 교사 최재석입니다. 오늘이 1정 연수 마지막주차라고 들었는데 얼마나 선생님들이 지친 상태일지 강사인 제가 걱정이 됩니다. 힘들었다고 대답하시더라구요. 제가 덜 힘들게 해드릴 순 없고, 편안하게 들으실 수 있게, 그리고 한가지라도 유익한 정보를 남길 수 있도록 최선에 다해보겠습니다.』"),
        # 2: Video 1
        ("p2", COLOR_BG_SLATE, "재생 전 한 문장만: 『디지털 성범죄하면 떠오르는 대표적인 사건은 2019년, 대학생 둘이 파고들어 밝혀낸 N번방 사건일 겁니다. 관련 영상 먼저 보고 이야기 나누도록 하겠습니다.』"),
        # 3: N-th Room 4 Steps
        ("p3", COLOR_BG_SLATE, "『N번방에 대해서 저도 당시에는 그렇게 관심이 많지 않았고, 생경하게 생각했었습니다. 알면알수록 분노가 치미는 이 사건에 대해서 자세히 알아보겠습니다. 일탈계 이용자를 대상으로 금전적 호의를 베풀고 개인정보를 기반으로 협박, 성착취로 이어지는 악질적 행위였습니다. 폭력이 아니라 개인정보를 쥐고 흔드는 협박입니다. 문형욱 34년, 조주빈 47년 4월 형 확정 받습니다.』"),
        # 4: Video 2
        ("p4", COLOR_BG_SLATE, "재생 전: 『그리고 5년 뒤. 또다른 디지털 성범죄가 학교를 뒤흔듭니다. 바로 ai를 활용한 딥페이크 범죄입니다. 관련 영상 한편 보도록 하겠습니다.』"),
        # 5: Deepfake Spread
        ("p5", COLOR_BG_SLATE, "『최근 우리에게 어려움을 주는 건 AI를 활용한 딥페이크 영상입니다. 몇몇 학교는 졸업앨범 교사 사진을 빼는 실정입니다. 수법은 간단합니다. \"재미있으니까 한번 해봐\"라며 받은 링크에 얼굴을 넣고 합성하고, 무료 회기가 끝나면 친구를 초대하게 만듭니다. 별 생각없이 사진을 넣고 음란물을 합성합니다. 장난삼아 했던 행동들이 범죄로 이어집니다.』"),
        # 6: Full Black
        ("p6", COLOR_BG_NAVY, "(2초 정적 후) 『뉴스 속 먼 나라 얘기가 아닙니다. 이미 우리 학교의, 그리고 교실의 문턱을 넘었습니다. 우리교과와는 아직 먼 이야기라고 생각이 드시나요? 오늘은 이런 위험에 대한 경고와 우리가 어떻게 대처해나가야 하는지를 생각해 봤으면 좋겠습니다.』"),
        # 7: Chapter I Cover
        ("p7", COLOR_BG_NAVY, "『방금 두 사건이, 특수교육 현장에 던지는 신호를 네 가지로 정리해보았습니다.』"),
        # 8: Warning 1 (1/4)
        ("p8", COLOR_BG_SLATE, "『가해자들의 친절한 접근·가스라이팅에 가장 취약한 게 우리 아이들일 수 있습니다. 인지적, 정서적 취약점을 파고드는 디지털 성범죄에 피해자가 언제든 될 수 있습니다.』"),
        # 9: Warning 2 (2/4)
        ("p9", COLOR_BG_SLATE, "『AI가 입힌 '놀이의 탈' — 가해의 인식 없이 가담하게 되는. 나도 모르게 가해자가 될 위험이 있습니다.』"),
        # 10: Warning 3 (3/4)
        ("p10", COLOR_BG_SLATE, "『여기서부턴 남 얘기가 아닙니다. 우리 모두 피해자가 될 수 있습니다. 일상·교육활동 사진이 노출되기 쉬운 특수교사 본인이 표적이 될 수 있습니다.』"),
        # 11: Warning 4 (4/4)
        ("p11", COLOR_BG_SLATE, "『비장애 중심 매뉴얼·사후 처리는 장애학생에게, 그리고 특수교육에 안 맞는 경향이 있습니다. 그렇기에 이런 성범죄 발생시 특수교육적 환경에 맞는 가이드라인이 필요합니다.』"),
        # 12: 3 Goals
        ("p12", COLOR_BG_SLATE, "『오늘 이렇게 세 가지 큰 주제를 목표로 삼아보았습니다. 위협을 알고, 발견·대응하고, 우리 아이 눈높이로 가르치기 위해 필요한 것들을 안내드리려고 합니다.』")
    ]

    for spec in slide_specs:
        pid, bg_col, script_notes = spec
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide, bg_col)
        add_notes(slide, script_notes)

        if pid == "p1":
            # Title Cover
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
            tf = tb.text_frame
            tf.word_wrap = True
            
            p0 = tf.paragraphs[0]
            p0.text = "1급 정교사 자격연수 · 특수교사 대상"
            p0.font.name = "Pretendard"
            p0.font.size = Pt(16)
            p0.font.bold = True
            p0.font.color.rgb = COLOR_GOLD

            p1 = tf.add_paragraph()
            p1.text = "더 쉽게 당하고,\n더 늦게 발견된다"
            p1.font.name = "Pretendard"
            p1.font.size = Pt(48)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_TEXT_MAIN

            p2 = tf.add_paragraph()
            p2.text = "\n장애학생 디지털 성범죄 예방과 교사의 역할"
            p2.font.name = "Pretendard"
            p2.font.size = Pt(22)
            p2.font.color.rgb = COLOR_TEXT_SUB

            p3 = tf.add_paragraph()
            p3.text = "\n강사 : 최재석 (노성중학교 교사)"
            p3.font.name = "Pretendard"
            p3.font.size = Pt(16)
            p3.font.color.rgb = COLOR_TEXT_SUB

        elif pid == "p2":
            add_header(slide, "PART 1. 오프닝 | 뉴스 영상 01")
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "[뉴스 영상] N번방, 그 시작"
            p.font.name = "Pretendard"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = COLOR_GOLD

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.5))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0, 0, 0)
            box.line.color.rgb = COLOR_GOLD
            
            p_box = box.text_frame.paragraphs[0]
            p_box.text = "▶ [영상 재생 영역: N번방_사건_요약.mp4]\n\n출처: YouTube (https://www.youtube.com/watch?v=VZgunHeAYvI)"
            p_box.font.name = "Pretendard"
            p_box.font.size = Pt(18)
            p_box.font.color.rgb = COLOR_TEXT_SUB
            p_box.alignment = PP_ALIGN.CENTER

        elif pid == "p3":
            add_header(slide, "PART 1. 오프닝 | 범죄 구조 분석")
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "N번방의 범죄 수법 4단계"
            p.font.name = "Pretendard"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_MAIN

            steps = [
                ("1. 친절한 제안", "금전/선물/호의 접근"),
                ("2. 개인정보 확보", "신분증/학교/얼굴 획득"),
                ("3. 협박 & 통제", "유포 협박 및 조종"),
                ("4. 성착취 강요", "착취물 제작/유포")
            ]
            for i, (st, sub) in enumerate(steps):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0 + i * 2.9), Inches(2.6), Inches(2.6), Inches(2.8))
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_CARD_BG
                card.line.color.rgb = COLOR_GOLD if i < 3 else COLOR_RED
                
                tf_c = card.text_frame
                tf_c.word_wrap = True
                p_c = tf_c.paragraphs[0]
                p_c.text = st
                p_c.font.name = "Pretendard"
                p_c.font.size = Pt(20)
                p_c.font.bold = True
                p_c.font.color.rgb = COLOR_GOLD if i < 3 else COLOR_RED
                
                p_sub = tf_c.add_paragraph()
                p_sub.text = f"\n{sub}"
                p_sub.font.name = "Pretendard"
                p_sub.font.size = Pt(14)
                p_sub.font.color.rgb = COLOR_TEXT_SUB

            # Stamp
            stamp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.8), Inches(4.8), Inches(0.9))
            stamp.fill.solid()
            stamp.fill.fore_color.rgb = RGBColor(40, 15, 20)
            stamp.line.color.rgb = COLOR_RED
            stamp.line.width = Pt(3)
            p_st = stamp.text_frame.paragraphs[0]
            p_st.text = "문형욱 34년 / 조주빈 47년 4개월 형 확정"
            p_st.font.name = "Pretendard"
            p_st.font.size = Pt(16)
            p_st.font.bold = True
            p_st.font.color.rgb = COLOR_RED
            p_st.alignment = PP_ALIGN.CENTER

        elif pid == "p4":
            add_header(slide, "PART 1. 오프닝 | 뉴스 영상 02")
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "[뉴스 영상] 2024, 딥페이크가 학교로"
            p.font.name = "Pretendard"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = COLOR_RED

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.5))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0, 0, 0)
            box.line.color.rgb = COLOR_RED
            
            p_box = box.text_frame.paragraphs[0]
            p_box.text = "▶ [영상 재생 영역: 학교_딥페이크_뉴스_요약.mp4]\n\n출처: YouTube (https://www.youtube.com/watch?v=FJU1OYdSJOY)"
            p_box.font.name = "Pretendard"
            p_box.font.size = Pt(18)
            p_box.font.color.rgb = COLOR_TEXT_SUB
            p_box.alignment = PP_ALIGN.CENTER

        elif pid == "p5":
            add_header(slide, "PART 1. 오프닝 | 딥페이크 메커니즘")
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "\"재밌으니 한번 해봐\" — 한 줄에서 7명으로 확산"
            p.font.name = "Pretendard"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = COLOR_GOLD

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.3), Inches(11.333), Inches(4.2))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_CARD_BG
            card.line.color.rgb = COLOR_GOLD
            
            tf_c = card.text_frame
            p_c = tf_c.paragraphs[0]
            p_c.text = "초기 1명 수신 ➔ 합성 텔레그램 봇 ➔ 친구 초대 강요가담 (7명 확산)\n\n"
            p_c.font.name = "Pretendard"
            p_c.font.size = Pt(24)
            p_c.font.bold = True
            p_c.font.color.rgb = COLOR_TEXT_MAIN
            
            p_c2 = tf_c.add_paragraph()
            p_c2.text = "⚠️ 현장 실태: 졸업앨범 교사 사진 전면 삭제 및 교사 대상 딥페이크 유포 사례 급증"
            p_c2.font.name = "Pretendard"
            p_c2.font.size = Pt(20)
            p_c2.font.bold = True
            p_c2.font.color.rgb = COLOR_RED

        elif pid == "p6":
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "디지털 성범죄,\n이미 교실 안에 있다"
            p.font.name = "Pretendard"
            p.font.size = Pt(52)
            p.font.bold = True
            p.font.color.rgb = COLOR_GOLD
            p.alignment = PP_ALIGN.CENTER

        elif pid == "p7":
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(3.5))
            tf = tb.text_frame
            p0 = tf.paragraphs[0]
            p0.text = "CHAPTER Ⅰ"
            p0.font.name = "JetBrains Mono"
            p0.font.size = Pt(24)
            p0.font.bold = True
            p0.font.color.rgb = COLOR_GOLD

            p1 = tf.add_paragraph()
            p1.text = "특수교육에 던지는 네 가지 경고"
            p1.font.name = "Pretendard"
            p1.font.size = Pt(44)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_TEXT_MAIN

        elif pid in ["p8", "p9", "p10", "p11"]:
            warn_idx = {"p8": "1 / 4", "p9": "2 / 4", "p10": "3 / 4", "p11": "4 / 4"}[pid]
            title_text = {
                "p8": "① 우리 학생도 피해자가 될 수 있다",
                "p9": "② 우리 학생도 가해자가 될 수 있다",
                "p10": "③ 교사도 예외는 아니다",
                "p11": "④ 맞춤형 가이드가 필요하다"
            }[pid]
            desc_text = {
                "p8": "인지적·정서적 취약성 집중 공략\n\n온라인에서의 거짓 호의와 가스라이팅 수법에 장애학생은 쉽게 경계심을 허물게 됩니다. '나를 인정해주는 친구'라는 착각을 유도하여 피해로 이끕니다.",
                "p9": "인식 없는 가담과 범죄의 장난화\n\nAI 딥페이크 봇이나 합성 게임은 단순한 '재미있는 놀이'로 인식됩니다. 이것이 심각한 성범죄이자 형사 처벌 대상임을 알지 못해 가담자로 전락합니다.",
                "p10": "교육활동 사진 노출과 딥페이크 위협\n\n학교 홈페이지, SNS, 졸업앨범 등에 사진이 자주 공개되는 교사 역시 딥페이크 범죄의 대상이 되고 있습니다. 더 이상 교사 안전 지대는 없습니다.",
                "p11": "비장애 중심 대응 매뉴얼의 한계\n\n복잡한 서면 진술 중심의 기존 매뉴얼은 표현이 어려운 장애학생에게 한계가 명확합니다. 장애 특성과 특수교육적 환경을 고려한 맞춤형 대응이 절실합니다."
            }[pid]
            
            is_red = (pid == "p10")
            accent_col = COLOR_RED if is_red else COLOR_GOLD

            add_header(slide, f"CHAPTER Ⅰ. 네 가지 경고 | PROGRESS {warn_idx}")
            
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = title_text
            p.font.name = "Pretendard"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = accent_col

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.4), Inches(11.333), Inches(4.2))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_CARD_BG
            card.line.color.rgb = accent_col
            card.line.width = Pt(2)
            
            tf_c = card.text_frame
            tf_c.word_wrap = True
            p_c = tf_c.paragraphs[0]
            p_c.text = desc_text
            p_c.font.name = "Pretendard"
            p_c.font.size = Pt(20)
            p_c.font.color.rgb = COLOR_TEXT_MAIN

        elif pid == "p12":
            add_header(slide, "CHAPTER Ⅰ. 마무리 | 학습 목표")
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "오늘, 세 가지 목표"
            p.font.name = "Pretendard"
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLOR_GOLD
            p.alignment = PP_ALIGN.CENTER

            goals = [
                ("GOAL A", "위협을 이해한다", "디지털 성범죄 본질 및\n장애학생 취약성 4구조 파악", COLOR_GOLD),
                ("GOAL B", "발견 · 대응한다", "4가지 조기 발견 신호 및\n실무 4-Track 대응 작동", COLOR_GOLD),
                ("GOAL C", "가르친다", "장애학생 눈높이에 맞춘\n경계·동의 및 예방 교육 적용", COLOR_TEAL)
            ]
            for i, (g_let, g_tit, g_sub, g_col) in enumerate(goals):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0 + i * 3.9), Inches(2.6), Inches(3.5), Inches(3.8))
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_CARD_BG
                card.line.color.rgb = g_col
                card.line.width = Pt(2)
                
                tf_c = card.text_frame
                tf_c.word_wrap = True
                
                p_let = tf_c.paragraphs[0]
                p_let.text = g_let
                p_let.font.name = "JetBrains Mono"
                p_let.font.size = Pt(22)
                p_let.font.bold = True
                p_let.font.color.rgb = g_col
                p_let.alignment = PP_ALIGN.CENTER

                p_tit = tf_c.add_paragraph()
                p_tit.text = f"\n{g_tit}"
                p_tit.font.name = "Pretendard"
                p_tit.font.size = Pt(24)
                p_tit.font.bold = True
                p_tit.font.color.rgb = COLOR_TEXT_MAIN
                p_tit.alignment = PP_ALIGN.CENTER

                p_sub = tf_c.add_paragraph()
                p_sub.text = f"\n\n{g_sub}"
                p_sub.font.name = "Pretendard"
                p_sub.font.size = Pt(16)
                p_sub.font.color.rgb = COLOR_TEXT_SUB
                p_sub.alignment = PP_ALIGN.CENTER

    output_path = r"g:\내 드라이브\2026 노성중\1급 정교사 연수 강의 요청\시각자료\강의_발표자료_챕터1.pptx"
    prs.save(output_path)
    print(f"v3 PPTX saved successfully to {output_path}")

if __name__ == "__main__":
    create_v3_chapter1_pptx()
